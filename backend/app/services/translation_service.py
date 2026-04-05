from typing import Optional
import os
import requests
import json
from PIL import Image
import pytesseract
import pdfplumber
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from dotenv import load_dotenv

# 加载环境变量
load_dotenv()

# 阿里云 API 配置
ALIYUN_API_KEY = os.getenv('ALIYUN_API_KEY', 'your-aliyun-api-key')
ALIYUN_MODEL_NAME = os.getenv('ALIYUN_MODEL_NAME', 'ep-20240405170950-xqxzh')

# 阿里云大模型 API 翻译函数
def aliyun_translate(text: str, target_lang: str, style: str = "general", industry: str = "general") -> str:
    """使用阿里云大模型 API 进行翻译"""
    try:
        # 构建请求 URL
        url = f"https://ark.cn-beijing.volces.com/api/v3/chat/completions"
        
        # 构建请求头
        headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {ALIYUN_API_KEY}"
        }
        
        # 构建请求体
        payload = {
            "model": ALIYUN_MODEL_NAME,
            "messages": [
                {
                    "role": "system",
                    "content": f"你是一个专业的翻译助手，需要将文本翻译成{target_lang}。翻译风格：{style}，行业规范：{industry}。"
                },
                {
                    "role": "user",
                    "content": f"请将以下文本翻译成{target_lang}：\n{text}"
                }
            ],
            "temperature": 0.7,
            "max_tokens": 2048
        }
        
        # 发送请求
        response = requests.post(url, headers=headers, data=json.dumps(payload))
        response.raise_for_status()
        
        # 解析响应
        result = response.json()
        translated_text = result['choices'][0]['message']['content']
        
        return translated_text
    except Exception as e:
        print(f"翻译失败: {str(e)}")
        return f"翻译失败: {str(e)}"

# 文字翻译
def translate_text(
    text: str,
    source_lang: str = "auto",
    target_lang: str = "zh",
    style: str = "general",
    industry: str = "general"
) -> str:
    """翻译文字"""
    # 使用阿里云大模型 API 进行翻译
    result = aliyun_translate(text, target_lang, style, industry)
    return result

# 图片翻译
def translate_image(image_path: str, target_lang: str = "zh") -> str:
    """翻译图片中的文字"""
    # 使用OCR识别图片中的文字
    try:
        img = Image.open(image_path)
        text = pytesseract.image_to_string(img, lang='eng+chi_sim')
        # 使用阿里云大模型 API 翻译识别出的文字
        result = aliyun_translate(text, target_lang)
        return result
    except Exception as e:
        return f"图片翻译失败: {str(e)}"

# 文件翻译
def translate_file(
    file_path: str,
    source_lang: str = "auto",
    target_lang: str = "zh",
    style: str = "general",
    industry: str = "general",
    translate_tables: bool = True,
    translate_images: bool = True
) -> str:
    """翻译文件"""
    file_ext = os.path.splitext(file_path)[1].lower()
    
    try:
        if file_ext == ".pdf":
            return translate_pdf(file_path, target_lang, style, industry, translate_tables, translate_images)
        elif file_ext == ".docx":
            return translate_docx(file_path, target_lang, style, industry, translate_tables, translate_images)
        elif file_ext == ".xlsx":
            return translate_xlsx(file_path, target_lang, style, industry, translate_tables)
        elif file_ext == ".pptx":
            return translate_pptx(file_path, target_lang, style, industry, translate_tables, translate_images)
        else:
            return f"不支持的文件格式: {file_ext}"
    except Exception as e:
        return f"文件翻译失败: {str(e)}"

# 翻译PDF文件
def translate_pdf(pdf_path: str, target_lang: str, style: str, industry: str, translate_tables: bool, translate_images: bool) -> str:
    """翻译PDF文件"""
    text = ""
    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            # 提取文本
            page_text = page.extract_text() or ""
            text += page_text
            
            # 提取表格
            if translate_tables:
                tables = page.extract_tables()
                for table in tables:
                    table_text = "\n".join(["\t".join([str(cell) if cell else "" for cell in row]) for row in table])
                    text += "\n" + table_text
    
    # 使用阿里云大模型 API 翻译提取的文本
    result = aliyun_translate(text, target_lang, style, industry)
    return result

# 翻译Word文件
def translate_docx(docx_path: str, target_lang: str, style: str, industry: str, translate_tables: bool, translate_images: bool) -> str:
    """翻译Word文件"""
    doc = Document(docx_path)
    text = ""
    
    # 提取段落
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    
    # 提取表格
    if translate_tables:
        for table in doc.tables:
            for row in table.rows:
                row_text = "\t".join([cell.text for cell in row.cells])
                text += row_text + "\n"
    
    # 使用阿里云大模型 API 翻译提取的文本
    result = aliyun_translate(text, target_lang, style, industry)
    return result

# 翻译Excel文件
def translate_xlsx(xlsx_path: str, target_lang: str, style: str, industry: str, translate_tables: bool) -> str:
    """翻译Excel文件"""
    if not translate_tables:
        return "未翻译表格内容"
    
    wb = load_workbook(xlsx_path)
    text = ""
    
    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        text += f"Sheet: {sheet_name}\n"
        
        for row in sheet.iter_rows(values_only=True):
            row_text = "\t".join([str(cell) if cell else "" for cell in row])
            text += row_text + "\n"
    
    # 使用阿里云大模型 API 翻译提取的文本
    result = aliyun_translate(text, target_lang, style, industry)
    return result

# 翻译PowerPoint文件
def translate_pptx(pptx_path: str, target_lang: str, style: str, industry: str, translate_tables: bool, translate_images: bool) -> str:
    """翻译PowerPoint文件"""
    prs = Presentation(pptx_path)
    text = ""
    
    for i, slide in enumerate(prs.slides):
        text += f"Slide {i+1}:\n"
        
        # 提取幻灯片中的文本
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
            
            # 提取表格
            if translate_tables and hasattr(shape, "table"):
                table = shape.table
                for row in table.rows:
                    row_text = "\t".join([cell.text for cell in row.cells])
                    text += row_text + "\n"
    
    # 使用阿里云大模型 API 翻译提取的文本
    result = aliyun_translate(text, target_lang, style, industry)
    return result